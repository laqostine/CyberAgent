"""Configuration settings for PowerPoint template layouts and formatting."""

from typing import Dict
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

# Define placeholder types using simple integers
# These values correspond to PowerPoint's internal placeholder types
TITLE = 1
CENTER_TITLE = 3
SUBTITLE = 4
BODY = 2
OBJECT = 7
DATE = 16
FOOTER = 15
SLIDE_NUMBER = 13
PICTURE = 18

# Template layout configuration
LAYOUT_MAPPING = {
    'title_slide': {
        'index': 0,
        'placeholders': {
            'title_1': {'type': CENTER_TITLE, 'idx': 0},
            'subtitle_2': {'type': SUBTITLE, 'idx': 1},
            'date_placeholder_3': {'type': DATE, 'idx': 10},
            'footer_placeholder_4': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_5': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'title_and_content': {
        'index': 1,
        'placeholders': {
            'title_1': {'type': TITLE, 'idx': 0},
            'content_placeholder_2': {'type': OBJECT, 'idx': 1},
            'date_placeholder_3': {'type': DATE, 'idx': 10},
            'footer_placeholder_4': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_5': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'section_header': {
        'index': 2,
        'placeholders': {
            'title_1': {'type': TITLE, 'idx': 0},
            'text_placeholder_2': {'type': BODY, 'idx': 1},
            'date_placeholder_3': {'type': DATE, 'idx': 10},
            'footer_placeholder_4': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_5': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'two_content': {
        'index': 3,
        'placeholders': {
            'title_1': {'type': TITLE, 'idx': 0},
            'content_placeholder_2': {'type': OBJECT, 'idx': 1},
            'content_placeholder_3': {'type': OBJECT, 'idx': 2},
            'date_placeholder_4': {'type': DATE, 'idx': 10},
            'footer_placeholder_5': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_6': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'comparison': {
        'index': 4,
        'placeholders': {
            'title_1': {'type': TITLE, 'idx': 0},
            'text_placeholder_2': {'type': BODY, 'idx': 1},
            'content_placeholder_3': {'type': OBJECT, 'idx': 2},
            'text_placeholder_4': {'type': BODY, 'idx': 3},
            'content_placeholder_5': {'type': OBJECT, 'idx': 4},
            'date_placeholder_6': {'type': DATE, 'idx': 10},
            'footer_placeholder_7': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_8': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'title_only': {
        'index': 5,
        'placeholders': {
            'title_1': {'type': TITLE, 'idx': 0},
            'date_placeholder_2': {'type': DATE, 'idx': 10},
            'footer_placeholder_3': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_4': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'blank': {
        'index': 6,
        'placeholders': {
            'date_placeholder_1': {'type': DATE, 'idx': 10},
            'footer_placeholder_2': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_3': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'content_with_caption': {
        'index': 7,
        'placeholders': {
            'title_1': {'type': TITLE, 'idx': 0},
            'content_placeholder_2': {'type': OBJECT, 'idx': 1},
            'text_placeholder_3': {'type': BODY, 'idx': 2},
            'date_placeholder_4': {'type': DATE, 'idx': 10},
            'footer_placeholder_5': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_6': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'picture_with_caption': {
        'index': 8,
        'placeholders': {
            'title_1': {'type': TITLE, 'idx': 0},
            'picture_placeholder_2': {'type': PICTURE, 'idx': 1},
            'text_placeholder_3': {'type': BODY, 'idx': 2},
            'date_placeholder_4': {'type': DATE, 'idx': 10},
            'footer_placeholder_5': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_6': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'title_and_vertical_text': {
        'index': 9,
        'placeholders': {
            'title_1': {'type': TITLE, 'idx': 0},
            'vertical_text_placeholder_2': {'type': BODY, 'idx': 1},
            'date_placeholder_3': {'type': DATE, 'idx': 10},
            'footer_placeholder_4': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_5': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
    'vertical_title_and_text': {
        'index': 10,
        'placeholders': {
            'vertical_title_1': {'type': TITLE, 'idx': 0},
            'vertical_text_placeholder_2': {'type': BODY, 'idx': 1},
            'date_placeholder_3': {'type': DATE, 'idx': 10},
            'footer_placeholder_4': {'type': FOOTER, 'idx': 11},
            'slide_number_placeholder_5': {'type': SLIDE_NUMBER, 'idx': 12},
        }
    },
}

# Slide type to layout mapping
SLIDE_TYPE_LAYOUTS = {
    'title': 'title_slide',
    'content': 'content_detailed',
    'section': 'section_content',
    'two_content': 'two_content'
}

# Layout configurations for different slide types
LAYOUT_CONFIG = {
    'content_detailed': {
        'index': 1,
        'placeholders': {
            'title': {'type': TITLE, 'idx': 0},
            'content': {'type': BODY, 'idx': 1}
        }
    },
    'section_content': {
        'index': 2,
        'placeholders': {
            'title': {'type': TITLE, 'idx': 0},
            'subtitle': {'type': SUBTITLE, 'idx': 1},
            'content': {'type': BODY, 'idx': 2}
        }
    }
}

# Text formatting settings for different element types
FORMATTING_CONFIG = {
    'main_title': {
        'font_size': 44,
        'font_bold': True,
        'font_name': 'Calibri',
        'alignment': PP_ALIGN.LEFT,
        'space_after': 12,
        'line_spacing': 1.2
    },
    'section_title': {
        'font_size': 36,
        'font_bold': True,
        'font_name': 'Calibri',
        'alignment': PP_ALIGN.LEFT,
        'space_after': 24,
        'line_spacing': 1.2
    },
    'body_large': {
        'font_size': 28,
        'font_bold': False,
        'font_name': 'Calibri',
        'line_spacing': 1.3,
        'space_after': 12,
        'alignment': PP_ALIGN.LEFT,
        'space_before': 6
    },
    'bullet_large': {
        'font_size': 24,
        'font_bold': False,
        'font_name': 'Calibri',
        'line_spacing': 1.2,
        'space_after': 6,
        'alignment': PP_ALIGN.LEFT,
        'indent_level': 1,
        'space_before': 3
    }
}

def get_layout_info(slide_type: str) -> dict:
    """Get layout configuration for a specific slide type."""
    return LAYOUT_CONFIG.get(slide_type, LAYOUT_CONFIG['content_detailed'])

def get_placeholder_info(layout_type: str, placeholder_name: str) -> dict:
    """Get placeholder information for a specific layout and placeholder name."""
    layout = LAYOUT_CONFIG.get(layout_type)
    if layout:
        return layout['placeholders'].get(placeholder_name)
    return None

def get_formatting(element_type: str) -> dict:
    """Get formatting settings for a specific element type."""
    return FORMATTING_CONFIG.get(element_type, FORMATTING_CONFIG['body_large'])

def format_long_text(text: str, max_width: int = 100) -> str:
    """Format long text with proper wrapping and indentation"""
    words = text.split()
    lines = []
    current_line = []
    current_width = 0
    
    for word in words:
        word_length = len(word)
        if current_width + word_length + 1 <= max_width:
            current_line.append(word)
            current_width += word_length + 1
        else:
            lines.append(' '.join(current_line))
            current_line = [word]
            current_width = word_length
            
    if current_line:
        lines.append(' '.join(current_line))
        
    return '\n'.join(lines)
