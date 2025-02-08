import os
import google.generativeai as genai
from PIL import Image
from dotenv import load_dotenv
import pytesseract
from pptx import Presentation
import io
from pptx.enum.shapes import MSO_SHAPE_TYPE
from copy import deepcopy

# Load environment variables
load_dotenv()

# Configure Gemini API
genai.configure(api_key=os.getenv('Gemini_API_KEY'))

# Initialize the model
model = genai.GenerativeModel('gemini-pro')

class GeminiContentEnhancer:
    def __init__(self):
        self.prompt_template = """You are a cybersecurity expert tasked with elaborating on brief descriptions of key data handling stages in cybersecurity.  You will receive text snippets, one at a time, focusing on either:

Data Collection
Data Cleaning and Preprocessing
Data Analysis
Data Visualization

For each text snippet provided, your objective is to generate a more detailed and comprehensive paragraph.  Enhance the original text by:

Adding Specificity: Incorporate concrete examples and scenarios relevant to cybersecurity.
Explaining Importance: Articulate the significance of the stage within the broader cybersecurity context.
Providing Context: Connect the stage to other cybersecurity processes and data handling workflows.
Elaborating on Techniques: Expand on the methods, tools, and best practices associated with each stage.
Maintaining Formal Tone: Ensure the output is informative, professional, and suitable for a cybersecurity audience.
Text you will provide must be under 100 words.

Please enhance the following text while maintaining its core message: """

    def extract_text_from_slide(self, slide):
        """Extract text from a PowerPoint slide"""
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_parts.append(text)
        return " ".join(text_parts)

    def enhance_content(self, text):
        """Send text to Gemini API and get enhanced content"""
        try:
            # Combine prompt template with input text
            full_prompt = f"{self.prompt_template}\n\n{text}"
            
            # Generate response from Gemini
            response = model.generate_content(full_prompt)
            
            return response.text
        except Exception as e:
            print(f"Error generating content: {e}")
            return None

    def process_presentation(self, pptx_path, output_path=None):
        """Process all slides in a PowerPoint presentation and update the content"""
        try:
            # Load the source presentation
            source_prs = Presentation(pptx_path)
            
            # Create output path with _enhanced suffix
            presentation_output = pptx_path.replace('.pptx', '_enhanced.pptx')
            
            # Create new presentation using the same file as template
            prs = Presentation(pptx_path)
            
            # Store current slides to remove after copying all content
            slides_to_remove = []
            for slide in prs.slides:
                slides_to_remove.append(slide)
            
            results = []
            
            # Process each slide
            for slide_number, source_slide in enumerate(source_prs.slides, 1):
                print(f"\nProcessing slide {slide_number}...")
                
                try:
                    # Find matching layout
                    source_layout = source_slide.slide_layout
                    matching_layout = None
                    
                    # Try to find the matching layout by name
                    for layout in prs.slide_masters[0].slide_layouts:
                        if layout.name == source_layout.name:
                            matching_layout = layout
                            break
                    
                    # If no match found by name, use the first layout
                    if matching_layout is None:
                        matching_layout = prs.slide_layouts[0]
                    
                    # Create new slide with matching layout
                    new_slide = prs.slides.add_slide(matching_layout)
                    
                    # Get title and content separately
                    title_text = ""
                    content_text = ""
                    
                    # Copy all shapes from source slide to new slide
                    for shape in source_slide.shapes:
                        if not hasattr(shape, "text"):
                            continue
                        
                        # Handle text content
                        if shape.is_placeholder:
                            if shape.placeholder_format.type == 1:  # Title
                                title_text = shape.text.strip()
                                # Copy title to new slide
                                if title_text and hasattr(new_slide, 'shapes'):
                                    title_shape = new_slide.shapes.title
                                    if title_shape:
                                        title_shape.text = title_text
                            else:  # Content
                                text = shape.text.strip()
                                if text and text != title_text:
                                    content_text += text + " "
                    
                    if content_text:
                        # Get enhanced content from Gemini
                        enhanced_content = self.enhance_content(content_text)
                        
                        if enhanced_content:
                            # Add enhanced content to new slide
                            for shape in new_slide.shapes:
                                if (hasattr(shape, "text") and 
                                    shape.is_placeholder and 
                                    shape.placeholder_format.type != 1):  # Not title
                                    shape.text = enhanced_content
                            
                            result = {
                                'slide_number': slide_number,
                                'title': title_text,
                                'original_text': content_text,
                                'enhanced_content': enhanced_content
                            }
                            results.append(result)
                            
                            # Print results
                            print(f"\nSlide {slide_number}:")
                            print("Title:", title_text)
                            print("Original Text:")
                            print(content_text)
                            print("\nEnhanced Content:")
                            print(enhanced_content)
                            print("-" * 80)
                    
                    # Remove original slides after processing all content
                    if slides_to_remove:
                        for slide in slides_to_remove:
                            rId = slide.slide_id
                            prs.part.drop_rel(rId)
                        slides_to_remove = []
                    
                    # Save after each slide is processed
                    prs.save(presentation_output)
                    print(f"Progress saved after slide {slide_number}")
                    
                except Exception as e:
                    print(f"Error processing slide {slide_number}: {e}")
                    # Save progress even if current slide fails
                    prs.save(presentation_output)
                    continue
            
            # Save text results if output path is provided
            if output_path:
                self.save_results(results, output_path)
            
            print(f"\nPresentation processing completed. Final version saved as: {presentation_output}")
            return results
            
        except Exception as e:
            print(f"Error processing presentation: {e}")
            return None

    def save_results(self, results, output_path):
        """Save results to a text file"""
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                for result in results:
                    f.write(f"\nSlide {result['slide_number']}:\n")
                    f.write(f"Title: {result['title']}\n")
                    f.write("Original Text:\n")
                    f.write(result['original_text'])
                    f.write("\n\nEnhanced Content:\n")
                    f.write(result['enhanced_content'])
                    f.write("\n" + "-" * 80 + "\n")
            print(f"\nResults saved to {output_path}")
        except Exception as e:
            print(f"Error saving results: {e}")

def main():
    # Example usage
    enhancer = GeminiContentEnhancer()
    
    # Replace with your PowerPoint file path (use raw string to handle backslashes)
    pptx_path = r"C:\Users\Bera\Desktop\Cyber agent\CyberAgent\Cyberagentsonsunumtest.pptx"
    output_path = "enhanced_content_results.txt"
    
    results = enhancer.process_presentation(pptx_path, output_path)

if __name__ == "__main__":
    main() 