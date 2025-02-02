from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
from typing import Dict
import os
import logging

logger = logging.getLogger(__name__)

# Define placeholder types manually
PLACEHOLDER_TYPES = {
    0: 'TITLE',
    1: 'BODY',
    2: 'CENTER_TITLE',
    3: 'SUBTITLE',
    4: 'TEXT',
    5: 'PICTURE',
    6: 'CHART',
    7: 'TABLE',
    8: 'FOOTER',
    9: 'HEADER',
    10: 'DATE',
    11: 'SLIDE_NUMBER',
    12: 'MEDIA',
    13: 'OBJECT'
}

def get_placeholder_type(type_id: int) -> str:
    """Convert placeholder type ID to string representation"""
    return PLACEHOLDER_TYPES.get(type_id, f'UNKNOWN_{type_id}')

def analyze_slide_layout(layout, idx: int) -> Dict:
    """Analyze a single slide layout"""
    layout_info = {
        'index': idx,
        'name': layout.name,
        'placeholders': []
    }
    
    # Analyze each placeholder in the layout
    for placeholder in layout.placeholders:
        ph_type = placeholder.placeholder_format.type
        ph_info = {
            'idx': placeholder.placeholder_format.idx,
            'type': get_placeholder_type(ph_type),
            'type_id': ph_type,
            'name': placeholder.name,
            'shape_type': str(placeholder.shape_type),
            'width': placeholder.width,
            'height': placeholder.height,
            'left': placeholder.left,
            'top': placeholder.top
        }
        layout_info['placeholders'].append(ph_info)
    
    return layout_info

def analyze_template(template_path: str, output_path: str = None) -> Dict:
    """Analyze PowerPoint template and return detailed information"""
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    prs = Presentation(template_path)
    template_info = {
        'slide_layouts': [],
        'slide_width': prs.slide_width,
        'slide_height': prs.slide_height
    }
    
    # Analyze each slide layout
    print("\n=== Template Analysis ===")
    print(f"Template file: {template_path}")
    print(f"Slide size: {prs.slide_width} x {prs.slide_height}")
    print("\nAnalyzing layouts...")
    
    for idx, layout in enumerate(prs.slide_layouts):
        layout_info = analyze_slide_layout(layout, idx)
        template_info['slide_layouts'].append(layout_info)
        
        # Print layout information
        print(f"\nLayout {idx}: {layout.name}")
        print("Placeholders:")
        for ph in layout_info['placeholders']:
            print(f"  - {ph['name']}")
            print(f"    Type: {ph['type']} (ID: {ph['type_id']})")
            print(f"    Index: {ph['idx']}")
            print(f"    Position: left={ph['left']}, top={ph['top']}")
            print(f"    Size: {ph['width']}x{ph['height']}")
    
    # Save to file if output path is provided
    if output_path:
        with open(output_path, 'w') as f:
            json.dump(template_info, f, indent=2)
        print(f"\nTemplate analysis saved to: {output_path}")
    
    return template_info

def create_placeholder_mapping(template_info: Dict) -> Dict:
    """Create a mapping of placeholder types and their locations"""
    mapping = {}
    
    for layout in template_info['slide_layouts']:
        layout_name = layout['name'].lower()
        layout_idx = layout['index']
        
        mapping[layout_idx] = {
            'name': layout['name'],
            'placeholders': {}
        }
        
        for ph in layout['placeholders']:
            mapping[layout_idx]['placeholders'][ph['idx']] = {
                'type': ph['type'],
                'type_id': ph['type_id'],
                'name': ph['name'],
                'position': {
                    'left': ph['left'],
                    'top': ph['top'],
                    'width': ph['width'],
                    'height': ph['height']
                }
            }
    
    return mapping

def generate_template_config(template_info: Dict, output_path: str = "template_config.py") -> None:
    """Generate a template configuration file based on the analysis"""
    layouts = {}
    for layout in template_info['slide_layouts']:
        layout_name = layout['name'].lower().replace(' ', '_')
        placeholders = {}
        for ph in layout['placeholders']:
            placeholders[ph['name'].lower().replace(' ', '_')] = {
                'type': ph['type_id'],
                'idx': ph['idx']
            }
        layouts[layout_name] = {
            'index': layout['index'],
            'placeholders': placeholders
        }
    
    # Generate the configuration file
    config_content = """# Auto-generated template configuration
from typing import Dict

# Template layout configuration
LAYOUT_MAPPING = {
"""
    
    for layout_name, layout_info in layouts.items():
        config_content += f"    '{layout_name}': {{\n"
        config_content += f"        'index': {layout_info['index']},\n"
        config_content += "        'placeholders': {\n"
        for ph_name, ph_info in layout_info['placeholders'].items():
            config_content += f"            '{ph_name}': {{'type': {ph_info['type']}, 'idx': {ph_info['idx']}}},\n"
        config_content += "        }\n    },\n"
    
    config_content += "}\n\n"
    config_content += """# Slide type to layout mapping
SLIDE_TYPE_LAYOUTS = {
    'title': 'title_slide',
    'content': 'content',
    'section': 'section',
    'two_content': 'two_content',
    'summary': 'content',
    'quiz': 'content',
    'lab': 'content'
}

def get_layout_info(layout_type: str) -> Dict:
    \"\"\"Get layout information for a specific type\"\"\"
    layout_name = SLIDE_TYPE_LAYOUTS.get(layout_type)
    return LAYOUT_MAPPING.get(layout_name) if layout_name else None

def get_placeholder_info(layout_type: str, placeholder_name: str) -> Dict:
    \"\"\"Get placeholder information for a specific layout and placeholder\"\"\"
    layout_info = get_layout_info(layout_type)
    if layout_info:
        return layout_info['placeholders'].get(placeholder_name)
    return None
"""
    
    with open(output_path, 'w') as f:
        f.write(config_content)
    print(f"\nTemplate configuration generated: {output_path}")

def main():
    # Analyze template
    template_path = "template.pptx"
    analysis_output = "template_analysis.json"
    config_output = "template_config.py"
    
    try:
        print("Starting template analysis...")
        template_info = analyze_template(template_path, analysis_output)
        
        # Create and display placeholder mapping
        print("\n=== Placeholder Mapping ===")
        mapping = create_placeholder_mapping(template_info)
        
        for layout_idx, layout_info in mapping.items():
            print(f"\nLayout {layout_idx}: {layout_info['name']}")
            print("Placeholders:")
            for ph_idx, ph_info in layout_info['placeholders'].items():
                print(f"  {ph_idx}: {ph_info['name']} ({ph_info['type']})")
        
        # Generate template configuration
        generate_template_config(template_info, config_output)
        
        print("\nAnalysis complete! You can now use the generated configuration files.")
        
    except Exception as e:
        print(f"Error analyzing template: {str(e)}")

if __name__ == "__main__":
    main() 