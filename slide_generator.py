from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from gemini_content_generator import GeminiContentGenerator, SlideContent, SlideType, LessonPlan
from template_config import (
    get_layout_info, 
    get_placeholder_info, 
    get_formatting,
    TITLE,
    CENTER_TITLE,
    SUBTITLE,
    BODY,
    OBJECT,
    DATE,
    FOOTER,
    SLIDE_NUMBER,
    PICTURE
)
import os
import logging
from dotenv import load_dotenv
from typing import Dict, List, Optional, Tuple

# Load environment variables
load_dotenv()
logger = logging.getLogger(__name__)

class EnhancedSlideGenerator:
    def __init__(self, template_path: str):
        """Initialize the slide generator with a template"""
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found: {template_path}")
        
        self.presentation = Presentation(template_path)
        self.template_path = template_path
        
        # Set default dimensions
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
        self.content_margin = Inches(0.5)
        self.title_height = Inches(1.5)
        
        # Content area dimensions
        self.content_width = self.slide_width - (2 * self.content_margin)
        self.content_height = self.slide_height - self.title_height - (2 * self.content_margin)
        
        # Initialize template analysis
        self._analyze_template()
        
        logger.info(f"Initialized slide generator with template: {template_path}")

    def _analyze_template(self):
        """Analyze template layouts and store metadata"""
        self.layout_info = {}
        for idx, layout in enumerate(self.presentation.slide_layouts):
            placeholders = {}
            for shape in layout.placeholders:
                placeholders[shape.placeholder_format.idx] = {
                    'type': shape.placeholder_format.type,
                    'name': shape.name,
                    'width': shape.width,
                    'height': shape.height,
                    'left': shape.left,
                    'top': shape.top
                }
            self.layout_info[idx] = {
                'name': layout.name,
                'placeholders': placeholders
            }
            logger.info(f"Layout {idx}: {len(placeholders)} placeholders - {[p['name'] for p in placeholders.values()]}")

    def _validate_content(self, content: SlideContent) -> bool:
        """Validate content structure and length"""
        try:
            # Check title length
            if len(content.title) > 100:
                logger.warning(f"Title too long ({len(content.title)} chars)")
                return False
            
            # Check main content
            if content.main_content:
                lines = content.main_content.split('\n')
                if len(lines) > 12:  # Increased to accommodate data analytics content
                    logger.warning(f"Main content too long ({len(lines)} lines)")
                    return False
                
                for line in lines:
                    if len(line) > 120:
                        logger.warning(f"Line too long ({len(line)} chars)")
                        return False
            
            # Check bullet points
            if content.bullet_points:
                if len(content.bullet_points) > 8:  # Increased for data analytics topics
                    logger.warning(f"Too many bullet points ({len(content.bullet_points)})")
                    return False
                
                for point in content.bullet_points:
                    if len(point) > 100:
                        logger.warning(f"Bullet point too long ({len(point)} chars)")
                        return False
            
            return True
            
        except Exception as e:
            logger.error(f"Content validation failed: {str(e)}")
            return False

    def _select_layout(self, content: SlideContent) -> int:
        """Select the most appropriate layout based on content"""
        try:
            # Get content characteristics
            has_bullets = bool(content.bullet_points)
            has_image = hasattr(content, 'image_path') and content.image_path
            content_length = len(content.main_content) if content.main_content else 0
            
            # Select layout based on content type and characteristics
            if content.slide_type == SlideType.TITLE:
                return 0  # Title slide
            elif has_image:
                return 8  # Picture with caption
            elif has_bullets and content_length > 500:
                return 3  # Two content layout
            elif has_bullets:
                return 1  # Title and content
            elif content_length > 1000:
                return 4  # Comparison layout
            else:
                return 1  # Default to title and content
                
        except Exception as e:
            logger.error(f"Layout selection failed: {str(e)}")
            return 1  # Default to title and content

    def _format_text_content(self, text: str, max_line_length: int = 80) -> str:
        """Format text content to fit within slide boundaries"""
        try:
            lines = []
            current_line = ""
            
            for word in text.split():
                if len(current_line) + len(word) + 1 <= max_line_length:
                    current_line += word + " "
                else:
                    lines.append(current_line.strip())
                    current_line = word + " "
            
            if current_line:
                lines.append(current_line.strip())
            
            return "\n".join(lines)
            
        except Exception as e:
            logger.error(f"Text formatting failed: {str(e)}")
            return text

    def _get_placeholder_by_type(self, slide, placeholder_type: int) -> Optional[object]:
        """Get a placeholder by its type"""
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == placeholder_type:
                return shape
        return None

    def _add_text_to_placeholder(self, placeholder, text: str, formatting: dict = None):
        """Add text to a placeholder with formatting"""
        if placeholder and hasattr(placeholder, 'text_frame'):
            text_frame = placeholder.text_frame
            text_frame.text = text
            
            if formatting:
                paragraph = text_frame.paragraphs[0]
                if 'font_size' in formatting:
                    paragraph.font.size = Pt(formatting['font_size'])
                if 'font_bold' in formatting:
                    paragraph.font.bold = formatting['font_bold']
                if 'font_name' in formatting:
                    paragraph.font.name = formatting['font_name']

    def create_slide_with_layout(self, layout_name: str) -> object:
        """Create a slide using a specific layout"""
        layout_info = get_layout_info(layout_name)
        if not layout_info:
            logger.warning(f"Layout '{layout_name}' not found, using default")
            return self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        
        return self.presentation.slides.add_slide(
            self.presentation.slide_layouts[layout_info['index']]
        )

    def add_content_to_slide(self, slide: object, content: SlideContent):
        """Add content to a slide based on its type"""
        try:
            # Find all placeholders in the slide
            placeholders = {}
            for shape in slide.placeholders:
                placeholders[shape.placeholder_format.type] = shape

            # Add title to title placeholder
            title_shape = placeholders.get(TITLE) or placeholders.get(CENTER_TITLE)
            if title_shape:
                title_frame = title_shape.text_frame
                title_frame.text = content.title
                self._apply_formatting(title_frame.paragraphs[0], get_formatting('main_title'))
                
                # If we have a subtitle in the content
                if hasattr(content, 'subtitle') and content.subtitle:
                    p = title_frame.add_paragraph()
                    p.text = content.subtitle
                    self._apply_formatting(p, get_formatting('section_title'))
            
            # Add main content to body placeholder
            body_shape = placeholders.get(BODY) or placeholders.get(OBJECT)
            if body_shape and content.main_content:
                text_frame = body_shape.text_frame
                text_frame.text = ""  # Clear default text
                text_frame.word_wrap = True
                
                # Format and add main content
                formatted_text = format_long_text(content.main_content, 80)  # Wrap at 80 chars
                p = text_frame.add_paragraph()
                p.text = formatted_text
                self._apply_formatting(p, get_formatting('body_large'))
                
                # Add bullet points with proper formatting
                if content.bullet_points:
                    for point in content.bullet_points:
                        p = text_frame.add_paragraph()
                        formatted_point = format_long_text(point, 60)  # Wrap bullet points at 60 chars
                        p.text = formatted_point
                        p.level = 1  # Indent level
                        self._apply_formatting(p, get_formatting('bullet_large'))
            
            # Only use fallback if no placeholders found
            if not (title_shape or body_shape):
                logger.warning("No placeholders found, using fallback method")
                self._add_fallback_content(slide, content)
            
        except Exception as e:
            logger.error(f"Error adding content to slide: {str(e)}")
            self._add_fallback_content(slide, content)

    def _apply_formatting(self, paragraph, formatting: dict):
        """Apply formatting to a paragraph"""
        if not paragraph or not formatting:
            return
            
        # Get the font object
        font = paragraph.font
        
        # Apply text formatting
        if 'font_size' in formatting:
            font.size = Pt(formatting['font_size'])
        if 'font_bold' in formatting:
            font.bold = formatting['font_bold']
        if 'font_name' in formatting:
            font.name = formatting['font_name']
        
        # Apply paragraph formatting
        if 'alignment' in formatting:
            paragraph.alignment = formatting['alignment']
        if 'line_spacing' in formatting:
            paragraph.line_spacing = formatting['line_spacing']
        if 'space_before' in formatting:
            paragraph.space_before = Pt(formatting['space_before'])
        if 'space_after' in formatting:
            paragraph.space_after = Pt(formatting['space_after'])
        if 'indent_level' in formatting:
            paragraph.level = formatting['indent_level']

    def _add_title_with_fallback(self, slide: object, title: str) -> bool:
        """Add title to slide with fallback mechanisms"""
        try:
            # Try primary title placeholder
            title_placeholder = self._get_placeholder_by_type(slide, TITLE)
            if title_placeholder:
                self._add_text_to_placeholder(title_placeholder, title, get_formatting('title'))
                return True
            
            # Try center title placeholder
            center_title = self._get_placeholder_by_type(slide, CENTER_TITLE)
            if center_title:
                self._add_text_to_placeholder(center_title, title, get_formatting('title'))
                return True
            
            # Fallback to textbox
            title_box = slide.shapes.add_textbox(
                self.content_margin,
                self.content_margin,
                self.slide_width - (2 * self.content_margin),
                self.title_height
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._apply_formatting(title_frame.paragraphs[0], get_formatting('title'))
            return True
            
        except Exception as e:
            logger.error(f"Error adding title: {str(e)}")
            return False

    def _add_main_content_with_fallback(self, slide: object, content: str) -> bool:
        """Add main content to slide with fallback mechanisms"""
        try:
            # Try body placeholder
            content_placeholder = self._get_placeholder_by_type(slide, BODY)
            if not content_placeholder:
                content_placeholder = self._get_placeholder_by_type(slide, OBJECT)
            
            if content_placeholder:
                self._add_text_to_placeholder(content_placeholder, content, get_formatting('body'))
                return True
            
            # Fallback to textbox
            content_box = slide.shapes.add_textbox(
                self.content_margin,
                self.title_height + self.content_margin,
                self.content_width,
                self.content_height
            )
            text_frame = content_box.text_frame
            text_frame.text = content
            text_frame.word_wrap = True
            self._apply_formatting(text_frame.paragraphs[0], get_formatting('body'))
            return True
            
        except Exception as e:
            logger.error(f"Error adding main content: {str(e)}")
            return False

    def _add_bullet_points_with_fallback(self, slide: object, bullet_points: List[str]) -> bool:
        """Add bullet points to slide with fallback mechanisms"""
        try:
            # Try body placeholder
            content_placeholder = self._get_placeholder_by_type(slide, BODY)
            if not content_placeholder:
                content_placeholder = self._get_placeholder_by_type(slide, OBJECT)
            
            if content_placeholder:
                text_frame = content_placeholder.text_frame
                text_frame.text = ""  # Clear default text
                
                for point in bullet_points:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    self._apply_formatting(p, get_formatting('bullet'))
                return True
            
            # Fallback to textbox
            content_box = slide.shapes.add_textbox(
                self.content_margin,
                self.title_height + self.content_margin,
                self.content_width,
                self.content_height
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for point in bullet_points:
                p = text_frame.add_paragraph()
                p.text = f"• {point}"
                p.level = 0
                self._apply_formatting(p, get_formatting('bullet'))
            return True
            
        except Exception as e:
            logger.error(f"Error adding bullet points: {str(e)}")
            return False

    def _add_fallback_content(self, slide: object, content: SlideContent):
        """Add content in a simple format when normal methods fail"""
        try:
            # Add title as simple textbox
            title_box = slide.shapes.add_textbox(
                self.content_margin,
                self.content_margin,
                self.slide_width - (2 * self.content_margin),
                self.title_height
            )
            title_frame = title_box.text_frame
            title_frame.text = content.title
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add content as simple textbox
            content_box = slide.shapes.add_textbox(
                self.content_margin,
                self.title_height + self.content_margin,
                self.content_width,
                self.content_height
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            # Add main content
            if content.main_content:
                p = text_frame.add_paragraph()
                p.text = content.main_content
            
            # Add bullet points
            if content.bullet_points:
                for point in content.bullet_points:
                    p = text_frame.add_paragraph()
                    p.text = f"• {point}"
                    p.level = 0
            
        except Exception as e:
            logger.error(f"Error adding fallback content: {str(e)}")
            # At this point, we can't do much more than log the error

    def generate_lesson_slides(self, lesson_plan: LessonPlan):
        """Generate all slides for a lesson"""
        try:
            # Create title slide
            title_slide = self.create_slide_with_layout('title')
            title_content = SlideContent(
                title="Data Analytics in Cybersecurity",
                main_content="Understanding Data Analytics and Its Applications in Cybersecurity",
                bullet_points=["Instructors: Ismail Molla, Ensar Bera"],
                slide_type=SlideType.TITLE
            )
            self.add_content_to_slide(title_slide, title_content)
            
            # Create course information slide
            info_slide = self.create_slide_with_layout('content')
            info_content = SlideContent(
                title="Course Information",
                main_content="Course Duration and Structure",
                bullet_points=[
                    "Lecture Duration: 4-8 hours",
                    "Lab Time: 1-2 hours",
                    "Interactive Sessions and Hands-on Practice",
                    "Real-world Case Studies and Examples"
                ],
                slide_type=SlideType.CONTENT
            )
            self.add_content_to_slide(info_slide, info_content)
            
            # Create section overview slide
            overview_slide = self.create_slide_with_layout('content')
            overview_content = SlideContent(
                title="Course Overview",
                main_content="Key Topics and Learning Objectives",
                bullet_points=[
                    "1. Data Collection - Understanding Sources and Methods",
                    "2. Data Cleaning and Preprocessing - Ensuring Data Quality",
                    "3. Data Analysis - Extracting Insights",
                    "4. Data Visualization - Presenting Results",
                    "5. Analytical Tools - Practical Implementation",
                    "6. Applications in Cybersecurity - Real-world Usage"
                ],
                slide_type=SlideType.CONTENT
            )
            self.add_content_to_slide(overview_slide, overview_content)
            
            # Create content slides for each section
            for i, slide_content in enumerate(lesson_plan.slides, 1):
                logger.info(f"Creating slide {i} of {len(lesson_plan.slides)}...")
                
                # Select appropriate layout based on content type
                layout_name = 'content'
                if slide_content.slide_type == SlideType.TITLE:
                    layout_name = 'section'
                elif hasattr(slide_content, 'image_path') and slide_content.image_path:
                    layout_name = 'picture'
                elif i % 5 == 0:  # Add summary slides periodically
                    layout_name = 'summary'
                
                slide = self.create_slide_with_layout(layout_name)
                self.add_content_to_slide(slide, slide_content)
                
                # Add practical examples for key concepts
                if i % 3 == 0:
                    example_slide = self.create_slide_with_layout('two_content')
                    example_content = self._create_example_content(slide_content.title)
                    self.add_content_to_slide(example_slide, example_content)
                
        except Exception as e:
            logger.error(f"Error generating lesson slides: {str(e)}")
            raise

    def _create_example_content(self, topic: str) -> SlideContent:
        """Create practical example content for a topic"""
        return SlideContent(
            title=f"Practical Example: {topic}",
            main_content="Real-world Application",
            bullet_points=[
                "Problem Statement",
                "Solution Approach",
                "Implementation Steps",
                "Results and Analysis"
            ],
            slide_type=SlideType.INTERACTIVE
        )

    def save_presentation(self, output_path: str):
        """Save the presentation to the specified path"""
        try:
            self.presentation.save(output_path)
            logger.info(f"Presentation saved successfully to {output_path}")
        except Exception as e:
            logger.error(f"Error saving presentation: {str(e)}")
            raise

    def add_visualization(self, slide, image_path: str, position: Optional[Tuple[float, float]] = None):
        """Add visualization (graph or icon) to slide"""
        try:
            if not os.path.exists(image_path):
                logger.error(f"Visualization file not found: {image_path}")
                return
                
            # Default position is center of content area if not specified
            if position is None:
                left = self.content_margin + (self.content_width / 4)  # Start at 1/4 of content width
                top = self.title_height + self.content_margin
            else:
                left = Inches(position[0])
                top = Inches(position[1])
            
            # Add image to slide
            pic = slide.shapes.add_picture(
                image_path,
                left,
                top,
                width=self.content_width / 2,  # Use half of content width by default
                height=self.content_height * 0.8  # Use 80% of content height
            )
            
            return pic
            
        except Exception as e:
            logger.error(f"Error adding visualization to slide: {str(e)}")
            return None

    def create_slide_with_visualization(self, title: str, content: str, visualization_path: Optional[str] = None):
        """Create a slide with both text content and visualization"""
        try:
            # Create slide
            slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])  # Blank layout
            
            # Add title
            title_box = slide.shapes.add_textbox(
                self.content_margin,
                self.content_margin,
                self.slide_width - (2 * self.content_margin),
                self.title_height
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            
            # Add content on the left side
            content_box = slide.shapes.add_textbox(
                self.content_margin,
                self.title_height + self.content_margin,
                self.content_width / 2,  # Use half width for content
                self.content_height
            )
            content_frame = content_box.text_frame
            content_frame.text = content
            content_frame.word_wrap = True
            
            # Add visualization on the right side if provided
            if visualization_path:
                self.add_visualization(
                    slide,
                    visualization_path,
                    (5.5, 2.0)  # Position on right side
                )
            
            return slide
            
        except Exception as e:
            logger.error(f"Error creating slide with visualization: {str(e)}")
            return None

def main():
    # Initialize Gemini content generator
    api_key = os.getenv("Gemini_API_KEY")
    if not api_key:
        raise ValueError("Please set Gemini_API_KEY in your .env file")
    
    content_generator = GeminiContentGenerator(api_key)
    slide_generator = EnhancedSlideGenerator("template.pptx")
    
    # Generate a complete lesson plan
    lesson_plan = content_generator.generate_lesson_plan(
        "data_collection",
        "Introduction to Data Sources"
    )
    
    # Generate all slides for the lesson
    slide_generator.generate_lesson_slides(lesson_plan)
    
    # Save the presentation
    output_path = "Updated_CyberAgent_Slide.pptx"
    slide_generator.save_presentation(output_path)
    print(f"Presentation saved as {output_path}")

if __name__ == "__main__":
    main() 