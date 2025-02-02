from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from gemini_content_generator import GeminiContentGenerator, SlideContent, SlideType, LessonPlan
import os
import logging
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

logger = logging.getLogger(__name__)

class SlideGenerator:
    def __init__(self, template_path: str):
        self.presentation = Presentation(template_path)
        self.template_path = template_path
        self._validate_template()

    def _validate_template(self):
        """Validate template layouts and log available placeholders"""
        logger.info("Validating template layouts...")
        self.layout_info = {}
        
        for idx, layout in enumerate(self.presentation.slide_layouts):
            placeholders = {}
            for shape in layout.placeholders:
                placeholders[shape.placeholder_format.idx] = {
                    'type': shape.placeholder_format.type,
                    'name': shape.name
                }
            self.layout_info[idx] = placeholders
            logger.info(f"Layout {idx}: {len(placeholders)} placeholders - {[p['name'] for p in placeholders.values()]}")

    def _get_safe_placeholder(self, slide, idx: int, fallback_method="add_textbox"):
        """Safely get a placeholder or create a textbox if not found"""
        try:
            return slide.shapes.placeholders[idx]
        except KeyError:
            logger.warning(f"Placeholder {idx} not found, using {fallback_method}")
            if fallback_method == "add_textbox":
                # Create a textbox instead
                left = Inches(1)
                top = Inches(2) if idx != 0 else Inches(1)  # Title higher up
                width = Inches(8)
                height = Inches(1.5)
                return slide.shapes.add_textbox(left, top, width, height)
            return None

    def create_title_slide(self, content: SlideContent) -> None:
        """Create a title slide with the given content"""
        slide = self.presentation.slides[0]
        
        # Title
        title_shape = self._get_safe_placeholder(slide, 0)
        if title_shape:
            title_shape.text = content.title
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(44)
                    run.font.bold = True
        
        # Subtitle/Description
        subtitle_shape = self._get_safe_placeholder(slide, 1)
        if subtitle_shape:
            subtitle_shape.text = content.main_content

    def create_content_slide(self, content: SlideContent) -> None:
        """Create a content slide based on the slide type"""
        layout_index = self._get_layout_index(content.slide_type)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[layout_index])
        
        # Add title
        title_shape = self._get_safe_placeholder(slide, 0)
        if title_shape:
            title_shape.text = content.title
        
        # Add content based on slide type
        try:
            if content.slide_type == SlideType.CONTENT:
                self._add_content_body(slide, content)
            elif content.slide_type == SlideType.INTERACTIVE:
                self._add_interactive_content(slide, content)
            elif content.slide_type == SlideType.LAB:
                self._add_lab_content(slide, content)
            elif content.slide_type == SlideType.QUIZ:
                self._add_quiz_content(slide, content)
            elif content.slide_type == SlideType.SUMMARY:
                self._add_summary_content(slide, content)
        except Exception as e:
            logger.error(f"Error adding content to slide: {str(e)}")
            self._add_fallback_content(slide, content)

    def _add_fallback_content(self, slide, content: SlideContent) -> None:
        """Add content in a simple format when normal methods fail"""
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = textbox.text_frame
        tf.word_wrap = True
        
        # Add main content
        p = tf.paragraphs[0]
        p.text = content.main_content
        
        # Add bullet points if any
        if content.bullet_points:
            for point in content.bullet_points:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.level = 0

    def _add_content_body(self, slide, content: SlideContent) -> None:
        """Add regular content with bullet points"""
        textbox = self._get_safe_placeholder(slide, 1)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        # Add main content
        tf.text = content.main_content
        
        # Add bullet points
        if content.bullet_points:
            for point in content.bullet_points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

    def _get_layout_index(self, slide_type: SlideType) -> int:
        """Get the appropriate layout index based on slide type"""
        layout_map = {
            SlideType.CONTENT: 1,
            SlideType.INTERACTIVE: 2,
            SlideType.LAB: 3,
            SlideType.QUIZ: 4,
            SlideType.SUMMARY: 5
        }
        return layout_map.get(slide_type, 1)  # Default to content layout

    def _add_interactive_content(self, slide, content: SlideContent) -> None:
        """Add interactive content with discussion points"""
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Discussion Points:"
        
        if content.interactive_elements:
            for point in content.interactive_elements.get("points", []):
                p = tf.add_paragraph()
                p.text = f"👉 {point}"
                p.level = 0

    def _add_lab_content(self, slide, content: SlideContent) -> None:
        """Add lab exercise content"""
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Lab Exercise"
        
        if content.interactive_elements:
            # Add objectives
            p = tf.add_paragraph()
            p.text = "Objectives:"
            for obj in content.interactive_elements.get("objectives", []):
                p = tf.add_paragraph()
                p.text = f"• {obj}"
                p.level = 1
            
            # Add tools
            p = tf.add_paragraph()
            p.text = "\nRequired Tools:"
            for tool in content.interactive_elements.get("tools", []):
                p = tf.add_paragraph()
                p.text = f"• {tool}"
                p.level = 1

    def _add_quiz_content(self, slide, content: SlideContent) -> None:
        """Add quiz questions"""
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Knowledge Check"
        
        if content.interactive_elements:
            for i, question in enumerate(content.interactive_elements.get("questions", []), 1):
                p = tf.add_paragraph()
                p.text = f"\nQ{i}: {question['question']}"
                for option in question.get("options", []):
                    p = tf.add_paragraph()
                    p.text = f"   □ {option}"
                    p.level = 1

    def _add_summary_content(self, slide, content: SlideContent) -> None:
        """Add summary content"""
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Key Takeaways:"
        
        if content.bullet_points:
            for point in content.bullet_points:
                p = tf.add_paragraph()
                p.text = f"✓ {point}"
                p.level = 0

    def generate_lesson_slides(self, lesson_plan: LessonPlan) -> None:
        """Generate all slides for a lesson"""
        try:
            # Create title slide
            logger.info("Creating title slide...")
            title_content = SlideContent(
                title=lesson_plan.title,
                main_content=lesson_plan.description,
                slide_type=SlideType.TITLE
            )
            self.create_title_slide(title_content)
            
            # Create content slides
            for i, slide in enumerate(lesson_plan.slides, 1):
                logger.info(f"Creating slide {i} of {len(lesson_plan.slides)}...")
                self.create_content_slide(slide)
                
        except Exception as e:
            logger.error(f"Error in generate_lesson_slides: {str(e)}", exc_info=True)
            raise

    def save_presentation(self, output_path: str) -> None:
        """Save the presentation to the specified path"""
        try:
            self.presentation.save(output_path)
            logger.info(f"Presentation saved successfully to {output_path}")
        except Exception as e:
            logger.error(f"Error saving presentation: {str(e)}")
            raise

def main():
    # Initialize Gemini content generator
    api_key = os.getenv("Gemini_API_KEY")  # Match the exact case from .env file
    if not api_key:
        raise ValueError("Please set Gemini_API_KEY in your .env file")
    
    content_generator = GeminiContentGenerator(api_key)
    slide_generator = SlideGenerator("template.pptx")
    
    # Generate a complete lesson plan
    lesson_plan = content_generator.generate_lesson_plan(
        "data_collection",
        "Introduction to Data Sources"
    )
    
    # Generate all slides for the lesson
    slide_generator.generate_lesson_slides(lesson_plan)
    
    # Save the presentation
    output_path = "Updated_CyberAgent_Slide.pptx"
    slide_generator.save_presentation(output_path)
    print(f"Presentation saved as {output_path}")

if __name__ == "__main__":
    main()
