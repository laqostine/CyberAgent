from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.placeholders import PP_PLACEHOLDER
from typing import Dict, Optional, List
import logging

logger = logging.getLogger(__name__)

class TemplateMapping:
    """Maps PowerPoint template placeholders and layouts"""
    
    # Standard PowerPoint placeholder types
    PLACEHOLDER_TYPES = {
        'title': PP_PLACEHOLDER.TITLE,
        'body': PP_PLACEHOLDER.BODY,
        'content': PP_PLACEHOLDER.OBJECT,
        'text': PP_PLACEHOLDER.TEXT,
        'subtitle': PP_PLACEHOLDER.SUBTITLE,
        'center_title': PP_PLACEHOLDER.CENTER_TITLE,
        'picture': PP_PLACEHOLDER.PICTURE,
        'chart': PP_PLACEHOLDER.CHART,
        'table': PP_PLACEHOLDER.TABLE,
        'footer': PP_PLACEHOLDER.FOOTER,
        'header': PP_PLACEHOLDER.HEADER,
        'date': PP_PLACEHOLDER.DATE,
        'slide_number': PP_PLACEHOLDER.SLIDE_NUMBER
    }

    # Layout types commonly used in templates
    LAYOUT_TYPES = {
        'title_slide': 0,        # Title slide layout
        'content': 1,            # Content slide layout
        'section_header': 2,     # Section header layout
        'two_content': 3,        # Two content layout
        'comparison': 4,         # Comparison layout
        'title_only': 5,         # Title only layout
        'blank': 6,             # Blank layout
        'content_caption': 7,    # Content with caption layout
        'picture_caption': 8     # Picture with caption layout
    }

    def __init__(self, template_path: str):
        """Initialize template mapping"""
        self.template_path = template_path
        self.presentation = Presentation(template_path)
        self.layout_map = {}
        self.placeholder_map = {}
        self._analyze_template()

    def _analyze_template(self) -> None:
        """Analyze template structure and map layouts and placeholders"""
        logger.info(f"Analyzing template: {self.template_path}")
        
        # Analyze each slide layout
        for idx, layout in enumerate(self.presentation.slide_layouts):
            layout_info = {
                'name': layout.name,
                'placeholders': {},
                'type': self._guess_layout_type(layout)
            }
            
            # Map placeholders in this layout
            for placeholder in layout.placeholders:
                ph_type = placeholder.placeholder_format.type
                ph_info = {
                    'idx': placeholder.placeholder_format.idx,
                    'type': ph_type,
                    'name': placeholder.name,
                    'shape_type': placeholder.shape_type
                }
                layout_info['placeholders'][ph_type] = ph_info
            
            self.layout_map[idx] = layout_info
            logger.info(f"Mapped layout {idx}: {layout.name} with {len(layout_info['placeholders'])} placeholders")

    def _guess_layout_type(self, layout) -> str:
        """Guess the layout type based on its characteristics"""
        name = layout.name.lower()
        
        if 'title slide' in name:
            return 'title_slide'
        elif 'section' in name:
            return 'section_header'
        elif 'two content' in name:
            return 'two_content'
        elif 'comparison' in name:
            return 'comparison'
        elif 'title only' in name:
            return 'title_only'
        elif 'blank' in name:
            return 'blank'
        elif 'picture' in name:
            return 'picture_caption'
        elif 'content' in name:
            return 'content'
        else:
            return 'unknown'

    def get_layout_by_type(self, layout_type: str) -> Optional[int]:
        """Get layout index by type"""
        for idx, info in self.layout_map.items():
            if info['type'] == layout_type:
                return idx
        return None

    def get_placeholder_by_type(self, layout_idx: int, placeholder_type: int) -> Optional[Dict]:
        """Get placeholder information by type"""
        if layout_idx in self.layout_map:
            placeholders = self.layout_map[layout_idx]['placeholders']
            return placeholders.get(placeholder_type)
        return None

    def get_title_placeholder(self, layout_idx: int) -> Optional[Dict]:
        """Get title placeholder for a layout"""
        return self.get_placeholder_by_type(layout_idx, PP_PLACEHOLDER.TITLE)

    def get_body_placeholder(self, layout_idx: int) -> Optional[Dict]:
        """Get body placeholder for a layout"""
        body = self.get_placeholder_by_type(layout_idx, PP_PLACEHOLDER.BODY)
        if not body:
            # Try object placeholder as fallback
            body = self.get_placeholder_by_type(layout_idx, PP_PLACEHOLDER.OBJECT)
        return body

    def get_layout_info(self, layout_idx: int) -> Optional[Dict]:
        """Get complete layout information"""
        return self.layout_map.get(layout_idx)

    def print_template_info(self) -> None:
        """Print detailed template information"""
        print("\n=== Template Analysis ===")
        for idx, layout_info in self.layout_map.items():
            print(f"\nLayout {idx}: {layout_info['name']} ({layout_info['type']})")
            print("Placeholders:")
            for ph_type, ph_info in layout_info['placeholders'].items():
                print(f"  - {ph_info['name']} (Type: {ph_type}, Index: {ph_info['idx']})")

def analyze_template(template_path: str) -> None:
    """Utility function to analyze a template file"""
    mapper = TemplateMapping(template_path)
    mapper.print_template_info()

if __name__ == "__main__":
    # Example usage
    template_path = "template.pptx"
    analyze_template(template_path) 